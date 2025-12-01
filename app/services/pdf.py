import fitz  # PyMuPDF
import httpx
import logging
import time
from urllib.parse import urlparse
from app.core.proxy import proxy_manager

logger = logging.getLogger(__name__)

async def download_and_extract(document_url: str) -> str:
    """
    Downloads a document (PDF, Word, PowerPoint) from the given URL and extracts text.
    Suporta: PDF, DOC, DOCX, PPT, PPTX.
    Registra métricas de tempo para download e extração.
    """
    start_ts = time.perf_counter()
    try:
        # Get a proxy to avoid blocking
        proxy = await proxy_manager.get_next_proxy()
        
        download_start = time.perf_counter()
        async with httpx.AsyncClient(proxy=proxy, verify=False) as client:
            # Add headers to mimic a browser, often needed for document downloads
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }
            response = await client.get(document_url, timeout=10.0, follow_redirects=True, headers=headers)
            response.raise_for_status()
            document_bytes = response.content
        download_duration = time.perf_counter() - download_start
        logger.info(
            f"[PERF] pdf step=download url={document_url} "
            f"duration={download_duration:.3f}s size_bytes={len(document_bytes)}"
        )

        # Determinar tipo de arquivo pela extensão
        parsed = urlparse(document_url)
        path_lower = parsed.path.lower()
        
        extract_start = time.perf_counter()
        if path_lower.endswith('.pdf'):
            text = await _extract_pdf_text(document_bytes, document_url)
            doc_type = "pdf"
        elif path_lower.endswith(('.doc', '.docx')):
            text = await _extract_word_text(document_bytes, document_url)
            doc_type = "word"
        elif path_lower.endswith(('.ppt', '.pptx')):
            text = await _extract_powerpoint_text(document_bytes, document_url)
            doc_type = "powerpoint"
        else:
            logger.warning(f"Tipo de documento não suportado: {document_url}")
            return ""

        extract_duration = time.perf_counter() - extract_start
        total_duration = time.perf_counter() - start_ts
        logger.info(
            f"[PERF] pdf step=extract url={document_url} type={doc_type} "
            f"duration={extract_duration:.3f}s total={total_duration:.3f}s chars={len(text)}"
        )
        return text

    except Exception as e:
        total_duration = time.perf_counter() - start_ts
        logger.warning(
            f"[PERF] pdf step=error url={document_url} total={total_duration:.3f}s error={e}"
        )
        return ""

async def _extract_pdf_text(pdf_bytes: bytes, pdf_url: str) -> str:
    """Extrai texto de PDF usando PyMuPDF"""
    extract_start = time.perf_counter()
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        text_parts = []
        total_pages = doc.page_count

        # Determine pages to extract: Only First 2 pages (0, 1) to save time and tokens
        indices = set()
        for i in range(min(2, total_pages)):
            indices.add(i)
        
        sorted_indices = sorted(list(indices))
        
        text_parts.append(f"--- PDF START: {pdf_url} ---")
        for i in sorted_indices:
            page = doc.load_page(i)
            text_parts.append(f"[Page {i+1}]\n{page.get_text()}")
        text_parts.append(f"--- PDF END: {pdf_url} ---")
        
        doc.close()
        text = "\n\n".join(text_parts)
        duration = time.perf_counter() - extract_start
        logger.debug(
            f"[PERF] pdf helper=_extract_pdf_text url={pdf_url} "
            f"duration={duration:.3f}s pages={len(sorted_indices)} chars={len(text)}"
        )
        return text
    except Exception as e:
        duration = time.perf_counter() - extract_start
        logger.warning(
            f"[PERF] pdf helper=_extract_pdf_text_error url={pdf_url} "
            f"duration={duration:.3f}s error={e}"
        )
        return ""

async def _extract_word_text(doc_bytes: bytes, doc_url: str) -> str:
    """Extrai texto de documentos Word (.doc, .docx)"""
    extract_start = time.perf_counter()
    try:
        # Tentar usar python-docx para .docx
        try:
            from docx import Document
            import io
            
            doc = Document(io.BytesIO(doc_bytes))
            text_parts = [f"--- WORD DOCUMENT START: {doc_url} ---"]
            
            # Extrair texto de parágrafos (primeiros 50 parágrafos para economizar tokens)
            for i, para in enumerate(doc.paragraphs[:50]):
                if para.text.strip():
                    text_parts.append(para.text)
            
            text_parts.append(f"--- WORD DOCUMENT END: {doc_url} ---")
            text = "\n\n".join(text_parts)
            duration = time.perf_counter() - extract_start
            logger.debug(
                f"[PERF] pdf helper=_extract_word_text url={doc_url} "
                f"duration={duration:.3f}s chars={len(text)}"
            )
            return text
        except ImportError:
            duration = time.perf_counter() - extract_start
            logger.warning(
                f"[PERF] pdf helper=_extract_word_text_import_error url={doc_url} "
                f"duration={duration:.3f}s error=python-docx-missing"
            )
            return f"--- WORD DOCUMENT FOUND: {doc_url} (processamento não disponível - instale python-docx) ---"
        except Exception as e:
            duration = time.perf_counter() - extract_start
            logger.warning(
                f"[PERF] pdf helper=_extract_word_text_error url={doc_url} "
                f"duration={duration:.3f}s error={e}"
            )
            return f"--- WORD DOCUMENT FOUND: {doc_url} (erro ao processar) ---"
    except Exception as e:
        duration = time.perf_counter() - extract_start
        logger.warning(
            f"[PERF] pdf helper=_extract_word_text_outer_error url={doc_url} "
            f"duration={duration:.3f}s error={e}"
        )
        return ""

async def _extract_powerpoint_text(ppt_bytes: bytes, ppt_url: str) -> str:
    """Extrai texto de apresentações PowerPoint (.ppt, .pptx)"""
    extract_start = time.perf_counter()
    try:
        # Tentar usar python-pptx para .pptx
        try:
            from pptx import Presentation
            import io
            
            prs = Presentation(io.BytesIO(ppt_bytes))
            text_parts = [f"--- POWERPOINT PRESENTATION START: {ppt_url} ---"]
            
            # Extrair texto de slides (primeiros 10 slides para economizar tokens)
            for i, slide in enumerate(prs.slides[:10]):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
            
            text_parts.append(f"--- POWERPOINT PRESENTATION END: {ppt_url} ---")
            text = "\n\n".join(text_parts)
            duration = time.perf_counter() - extract_start
            logger.debug(
                f"[PERF] pdf helper=_extract_powerpoint_text url={ppt_url} "
                f"duration={duration:.3f}s chars={len(text)}"
            )
            return text
        except ImportError:
            duration = time.perf_counter() - extract_start
            logger.warning(
                f"[PERF] pdf helper=_extract_powerpoint_text_import_error url={ppt_url} "
                f"duration={duration:.3f}s error=python-pptx-missing"
            )
            return f"--- POWERPOINT PRESENTATION FOUND: {ppt_url} (processamento não disponível - instale python-pptx) ---"
        except Exception as e:
            duration = time.perf_counter() - extract_start
            logger.warning(
                f"[PERF] pdf helper=_extract_powerpoint_text_error url={ppt_url} "
                f"duration={duration:.3f}s error={e}"
            )
            return f"--- POWERPOINT PRESENTATION FOUND: {ppt_url} (erro ao processar) ---"
    except Exception as e:
        duration = time.perf_counter() - extract_start
        logger.warning(
            f"[PERF] pdf helper=_extract_powerpoint_text_outer_error url={ppt_url} "
            f"duration={duration:.3f}s error={e}"
        )
        return ""

